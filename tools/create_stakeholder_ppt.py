from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "Chubb_LegalEntity_Change_Stakeholder_Deck.pptx"

NAVY = RGBColor(18, 42, 66)
TEAL = RGBColor(0, 133, 143)
GREEN = RGBColor(38, 126, 85)
AMBER = RGBColor(224, 153, 37)
RED = RGBColor(176, 62, 62)
PALE = RGBColor(239, 246, 248)
DARK = RGBColor(45, 55, 65)
WHITE = RGBColor(255, 255, 255)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def textbox(slide, text, x, y, w, h, size=18, color=DARK, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align is not None:
        p.alignment = align
    return box


def title(slide, heading, subtitle=None):
    textbox(slide, heading, 0.55, 0.3, 12.2, 0.5, 26, NAVY, True)
    if subtitle:
        textbox(slide, subtitle, 0.58, 0.86, 12.0, 0.35, 11, TEAL)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(1.22), Inches(1.3), Inches(0.06))
    line.fill.solid(); line.fill.fore_color.rgb = TEAL; line.line.fill.background()


def footer(slide, number):
    textbox(slide, "Chubb Enterprise Ontology Governance | ONTO-453 | Approver: Rajesh Gupta", 0.55, 7.16, 11.8, 0.2, 8, RGBColor(100, 110, 120))
    textbox(slide, str(number), 12.55, 7.13, 0.3, 0.22, 9, RGBColor(100, 110, 120), False, PP_ALIGN.RIGHT)


def bullet_list(slide, items, x, y, w, h, size=18, color=DARK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.level = 0; p.font.name = "Aptos"; p.font.size = Pt(size); p.font.color.rgb = color
        p.space_after = Pt(12)
        p.text = "• " + p.text
    return box


def card(slide, label, body, x, y, w, h, color=PALE, accent=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.color.rgb = accent
    textbox(slide, label, x + 0.18, y + 0.14, w - 0.36, 0.35, 15, accent, True)
    textbox(slide, body, x + 0.18, y + 0.58, w - 0.36, h - 0.7, 12, DARK)

# 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background.fill; background.solid(); background.fore_color.rgb = NAVY
textbox(slide, "Foundation Ontology Change", 0.75, 1.45, 11.8, 0.7, 34, WHITE, True)
textbox(slide, "Party → LegalEntity", 0.78, 2.3, 11.5, 0.65, 30, RGBColor(116, 221, 215), True)
textbox(slide, "Stakeholder operating procedure and governance evidence", 0.8, 3.2, 10.8, 0.4, 18, WHITE)
textbox(slide, "Ticket ONTO-453  |  Release 2.0.0  |  Approved by Rajesh Gupta", 0.8, 5.9, 11.5, 0.3, 13, RGBColor(210, 225, 235))

# 2
slide = prs.slides.add_slide(prs.slide_layouts[6]); title(slide, "1. Change requested", "What changed in the foundation ontology")
card(slide, "Before | release 1.0.0", "core:Party was the parent class for core:Person and core:Organization.", 0.7, 1.7, 5.7, 2.0, RGBColor(250, 242, 242), RED)
card(slide, "After | release 2.0.0", "core:LegalEntity replaces core:Party; Person and Organization now subclass LegalEntity.", 6.95, 1.7, 5.7, 2.0, RGBColor(237, 248, 242), GREEN)
textbox(slide, "The class IRI changes from .../core/Party to .../core/LegalEntity. This changes inferred meaning for downstream consumers.", 1.0, 4.35, 11.2, 0.9, 20, NAVY, True, PP_ALIGN.CENTER)
footer(slide, 2)

# 3
slide = prs.slides.add_slide(prs.slide_layouts[6]); title(slide, "2. Developer execution flow", "Run the same checks locally before requesting approval")
steps = [
    ("Edit", "ontology/fnd/core.ttl"),
    ("Detect", "python scripts/watch.py"),
    ("Compare", "RDF semantic diff"),
    ("Validate", "SHACL + consistency + tests"),
    ("Classify", "MAJOR → 2.0.0"),
]
for i, (label, body) in enumerate(steps):
    x = 0.55 + i * 2.55
    card(slide, label, body, x, 2.0, 2.2, 1.55, PALE, TEAL if i < 4 else AMBER)
    if i < len(steps) - 1:
        textbox(slide, "→", x + 2.24, 2.5, 0.3, 0.3, 24, TEAL, True, PP_ALIGN.CENTER)
bullet_list(slide, ["Baseline command: python scripts/govern.py --ci", "Expected result: Validation PASS, semantic impact MAJOR, suggested version 2.0.0"], 1.0, 4.45, 11.0, 1.2, 17)
footer(slide, 3)

# 4
slide = prs.slides.add_slide(prs.slide_layouts[6]); title(slide, "3. Governance decision", "Why this is a MAJOR release")
card(slide, "Removed", "core:Party entity and its subclass axioms", 0.7, 1.7, 3.8, 1.55, RGBColor(250, 242, 242), RED)
card(slide, "Added", "core:LegalEntity entity and hierarchy", 4.78, 1.7, 3.8, 1.55, RGBColor(237, 248, 242), GREEN)
card(slide, "Impact", "Consumers querying Party may break or return different inferred results", 8.86, 1.7, 3.8, 1.55, RGBColor(255, 248, 230), AMBER)
bullet_list(slide, ["MAJOR means approval is required before release.", "The release process preserves an immutable 1.0.0 snapshot for rollback and audit.", "Downstream mappings, queries, APIs, and AI/RAG consumers must adopt LegalEntity."], 1.0, 4.1, 11.2, 1.8, 17)
footer(slide, 4)

# 5
slide = prs.slides.add_slide(prs.slide_layouts[6]); title(slide, "4. Approval and release procedure", "Approved execution completed for ONTO-453")
steps = [
    ("1", "Review report", "reports/latest/change-report.md"),
    ("2", "Approve", "Rajesh Gupta"),
    ("3", "Release", "scripts/release.py"),
    ("4", "Publish", "Git commit + push"),
]
for i, (num, label, body) in enumerate(steps):
    y = 1.65 + i * 1.18
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), Inches(y), Inches(0.52), Inches(0.52))
    circ.fill.solid(); circ.fill.fore_color.rgb = TEAL; circ.line.fill.background()
    textbox(slide, num, 0.9, y + 0.1, 0.52, 0.25, 14, WHITE, True, PP_ALIGN.CENTER)
    textbox(slide, label, 1.7, y, 2.2, 0.3, 17, NAVY, True)
    textbox(slide, body, 4.0, y, 7.8, 0.35, 16, DARK)
textbox(slide, 'python scripts/release.py --approved-by "Rajesh Gupta" --ticket "ONTO-453" --reason "Replace Party with LegalEntity in the foundation ontology"', 1.0, 6.0, 11.2, 0.55, 11, TEAL, True)
footer(slide, 5)

# 6
slide = prs.slides.add_slide(prs.slide_layouts[6]); title(slide, "5. Evidence produced", "Traceability for reviewers and consumers")
items = [
    ("Release snapshot", "releases/2.0.0/"),
    ("Manifest", "releases/2.0.0/manifest.json"),
    ("Change record", "provenance/CHG-2026-3472F7.ttl"),
    ("Release event", "events/ontology-release-latest.json"),
    ("Registry state", "governance/release-state.json"),
    ("Governance report", "reports/latest/change-report.json and .md"),
]
for i, (label, body) in enumerate(items):
    x = 0.8 + (i % 2) * 6.1; y = 1.6 + (i // 2) * 1.35
    card(slide, label, body, x, y, 5.55, 1.0, PALE, TEAL)
footer(slide, 6)

# 7
slide = prs.slides.add_slide(prs.slide_layouts[6]); title(slide, "6. Consumer migration and rollback", "Actions after stakeholder approval")
card(slide, "Migrate", "Update SPARQL, mappings, APIs, documentation, and AI/RAG prompts from Party to LegalEntity. Re-run dependent tests.", 0.7, 1.7, 5.7, 2.25, RGBColor(237, 248, 242), GREEN)
card(slide, "Rollback", "Keep releases/1.0.0 immutable. If validation or consumer checks fail, stop publication and restore the approved release baseline with scripts/demo_change.py reset.", 6.95, 1.7, 5.7, 2.25, RGBColor(250, 242, 242), RED)
bullet_list(slide, ["Do not edit releases/1.0.0 or releases/2.0.0.", "Any follow-up ontology change must go through the same semantic diff and approval gate.", "Release event consumers refresh only from approved release events."], 1.0, 4.65, 11.0, 1.3, 17)
footer(slide, 7)

# 8
slide = prs.slides.add_slide(prs.slide_layouts[6]); title(slide, "Stakeholder sign-off", "Decision requested")
textbox(slide, "Approve the LegalEntity vocabulary and hierarchy change for release 2.0.0.", 1.0, 1.7, 11.2, 0.65, 25, NAVY, True, PP_ALIGN.CENTER)
card(slide, "Governance result", "PASS | MAJOR | 2.0.0", 1.3, 3.0, 3.2, 1.45, RGBColor(237, 248, 242), GREEN)
card(slide, "Approver", "Rajesh Gupta", 5.05, 3.0, 3.2, 1.45, PALE, TEAL)
card(slide, "Ticket", "ONTO-453", 8.8, 3.0, 3.2, 1.45, RGBColor(255, 248, 230), AMBER)
textbox(slide, "Next step: coordinate consumer migration and monitor release event adoption.", 1.0, 5.35, 11.2, 0.5, 19, DARK, True, PP_ALIGN.CENTER)
footer(slide, 8)

prs.save(OUT)
print(OUT)
